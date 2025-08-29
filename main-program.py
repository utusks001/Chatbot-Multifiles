# app.py
import os
from io import BytesIO
import requests
import streamlit as st
from dotenv import load_dotenv

# File parsing (existing behavior kept)
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PIL import Image

# Data analysis
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

# LangChain / VectorStore / Embeddings / LLM (existing)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# -------------------------
# Config / env
# -------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY")

st.set_page_config(
    page_title="Gemini + Groq Multi-file Chatbot (FAISS + OCR.Space) + Excel Analysis",
    page_icon="🤖",
    layout="wide"
)

if not (GOOGLE_API_KEY or GROQ_API_KEY):
    st.error("❌ GOOGLE_API_KEY atau GROQ_API_KEY tidak ditemukan. Tambahkan ke file .env sebelum menjalankan.")
    st.stop()

# -------------------------
# Embeddings & splitter (existing)
# -------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# -------------------------
# Session state initialization (ensure consistent keys)
# -------------------------
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []
if "dataframes" not in st.session_state:
    # structure: { filename: { "sheets": { sheet_name: df, ... } }, ... }
    st.session_state.dataframes = {}

# -------------------------
# Helpers for DataFrame analysis
# -------------------------
def safe_describe(df):
    """Describe that works on older pandas versions too"""
    try:
        return df.describe(include="all", datetime_is_numeric=True)
    except TypeError:
        return df.describe(include="all")

def df_info_text(df):
    import io
    buf = io.StringIO()
    df.info(buf=buf)
    return buf.getvalue()

def df_to_index_text(df, filename, sheet_name):
    """Create a short textual summary of dataframe suitable for indexing."""
    rows, cols = df.shape
    stats = safe_describe(df).transpose().reset_index().to_string(index=False)
    sample = df.head(20).to_csv(index=False)
    return f"DATAFRAME — file={filename}, sheet={sheet_name}\nshape: {rows}x{cols}\n{stats}\nSAMPLE:\n{sample}"

# -------------------------
# Existing file extractors (keep behavior for non-tabular unchanged)
# -------------------------
def extract_text_from_pdf(file_bytes: BytesIO):
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text

def extract_text_from_txt(file_bytes: BytesIO):
    try:
        file_bytes.seek(0)
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""

def extract_text_from_docx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text

def extract_text_from_pptx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text

# -------------------------
# OCR.Space Extractor (Image Files) — fixed to avoid zero-size sends
# -------------------------
def extract_text_from_image(file_bytes: BytesIO, filename="upload.png"):
    if not OCR_SPACE_API_KEY:
        st.warning("⚠️ OCR_SPACE_API_KEY tidak ditemukan di .env — image OCR dinonaktifkan.")
        return ""
    try:
        # ensure pointer at start and read bytes to validate
        file_bytes.seek(0)
        data = file_bytes.read()
        if not data or len(data) == 0:
            st.warning(f"⚠️ Gambar {filename} kosong atau gagal terbaca.")
            return ""
        # send BytesIO(data) so requests sends fresh stream
        resp = requests.post(
            "https://api.ocr.space/parse/image",
            files={"file": (filename, BytesIO(data), "image/png")},
            data={"apikey": OCR_SPACE_API_KEY, "language": "eng"},
            timeout=60
        )
        result = resp.json()
        if result.get("IsErroredOnProcessing"):
            st.warning("⚠️ OCR.Space error: " + str(result.get("ErrorMessage", ["Unknown error"])))
            return ""
        parsed = [p.get("ParsedText", "") for p in result.get("ParsedResults", [])]
        return "\n".join(parsed).strip()
    except Exception as e:
        st.warning(f"⚠️ OCR error: {e}")
        return ""

# -------------------------
# CSV / Excel extractors (NEW)
#   - read CSV / XLS / XLSX into pandas DataFrame
#   - store DataFrame into st.session_state.dataframes
#   - return a short text summary for indexing into FAISS
# -------------------------
def extract_text_from_csv(file_bytes: BytesIO, filename: str):
    try:
        file_bytes.seek(0)
        # try default, fallback to latin-1
        try:
            df = pd.read_csv(file_bytes)
        except Exception:
            file_bytes.seek(0)
            df = pd.read_csv(file_bytes, encoding="latin-1")
        st.session_state.dataframes[filename] = {"sheets": {"CSV": df}}
        return df_to_index_text(df, filename, "CSV")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca CSV {filename}: {e}")
        return ""

def extract_text_from_excel(file_bytes: BytesIO, filename: str):
    text_parts = []
    try:
        file_bytes.seek(0)
        xls = pd.ExcelFile(file_bytes)
        sheet_map = {}
        for s in xls.sheet_names:
            try:
                df = xls.parse(s)
                sheet_map[s] = df
                text_parts.append(df_to_index_text(df, filename, s))
            except Exception as se:
                st.warning(f"⚠️ Gagal parse sheet '{s}' di {filename}: {se}")
        if sheet_map:
            st.session_state.dataframes[filename] = {"sheets": sheet_map}
        return "\n\n---\n\n".join(text_parts)
    except Exception as e:
        st.warning(f"⚠️ Gagal baca Excel {filename}: {e}")
        return ""

# -------------------------
# Generic extractor (integrates all)
# -------------------------
def extract_text_from_file(uploaded_file):
    name = uploaded_file.name
    lname = name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    if lname.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    if lname.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    if lname.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    if lname.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    if lname.endswith(".csv"):
        # new CSV handling
        return extract_text_from_csv(BytesIO(raw), name)
    if lname.endswith((".xlsx", ".xls")):
        # new Excel handling
        return extract_text_from_excel(BytesIO(raw), name)
    if lname.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".jfif")):
        return extract_text_from_image(BytesIO(raw), filename=name)
    # keep original warnings for older formats
    if lname.endswith(".doc") or lname.endswith(".ppt"):
        st.warning(f"⚠️ File `{uploaded_file.name}` berformat lama (.doc/.ppt). Silakan konversi ke .docx/.pptx.")
        return ""
    st.warning(f"⚠️ Tipe file `{uploaded_file.name}` tidak didukung.")
    return ""

# -------------------------
# Build documents & FAISS (existing behavior)
# -------------------------
def build_documents_from_uploads(uploaded_files):
    docs = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs

def build_faiss_from_documents(docs):
    if not docs:
        return None
    vs = FAISS.from_documents(docs, embedding=EMBEDDINGS)
    return vs

# -------------------------
# Analysis UI for DataFrames (NEW)
#   - head(10), tail(10), describe(), info()
#   - outlier boxplots for Sales, Quantity, Profit (if present)
#   - scatter plot Sales vs Profit (if present)
#   - correlation heatmap (numeric)
#   - flexible Top/Bottom-N
#   - export Excel and HTML
# -------------------------
def auto_analyze_dataframe(df: pd.DataFrame, filename: str, sheet_name: str, show_in_app: bool = True):
    num_df = df.select_dtypes(include="number")

    # Prepare Excel export
    out_excel = BytesIO()
    with pd.ExcelWriter(out_excel, engine="openpyxl") as writer:
        try:
            df.to_excel(writer, sheet_name="Data", index=False)
        except Exception:
            df.reset_index().to_excel(writer, sheet_name="Data", index=False)
        # describe
        try:
            safe_describe(df).to_excel(writer, sheet_name="Describe")
        except Exception:
            pass
        if not num_df.empty:
            num_df.corr().to_excel(writer, sheet_name="Correlation")
    out_excel.seek(0)

    # Prepare HTML export (simple)
    html_body = f"<h2>Analysis — {filename} / {sheet_name}</h2><pre>{df.head(20).to_string()}</pre>"
    html_bytes = html_body.encode("utf-8")

    if show_in_app:
        st.markdown(f"### 📄 Analisa: {filename} — {sheet_name}")

        # Head & Tail
        st.write("**Head (10):**")
        st.dataframe(df.head(10))
        st.write("**Tail (10):**")
        st.dataframe(df.tail(10))

        # describe & info
        st.write("**describe():**")
        try:
            st.dataframe(safe_describe(df))
        except Exception as e:
            st.text(f"describe() error: {e}")

        st.write("**info():**")
        st.text(df_info_text(df))

        # Outlier Detection (Boxplot) only for Sales, Quantity, Profit (if present)
        target_cols = [c for c in ["Sales", "Quantity", "Profit"] if c in df.columns]
        if target_cols:
            st.write("**Outlier Detection (Boxplot)**")
            # show each boxplot (small figures)
            for col in target_cols:
                fig, ax = plt.subplots(figsize=(5, 3))
                # use seaborn for nicer boxplot; if seaborn missing, fallback to matplotlib
                try:
                    sns.boxplot(x=df[col].dropna(), ax=ax)
                except Exception:
                    ax.boxplot(df[col].dropna(), vert=False)
                ax.set_title(f"Outliers — {col}")
                st.pyplot(fig)

        # Scatter plot Sales vs Profit (if both present)
        if "Sales" in df.columns and "Profit" in df.columns:
            st.write("**Scatter Plot: Sales vs Profit**")
            fig, ax = plt.subplots(figsize=(5, 3))
            try:
                sns.scatterplot(x=df["Sales"], y=df["Profit"], ax=ax)
            except Exception:
                ax.scatter(df["Sales"], df["Profit"], s=10)
            ax.set_xlabel("Sales")
            ax.set_ylabel("Profit")
            st.pyplot(fig)

        # Correlation heatmap for numeric columns
        if not num_df.empty:
            st.write("**Correlation matrix:**")
            corr = num_df.corr()
            st.dataframe(corr)
            fig, ax = plt.subplots(figsize=(5, 3))
            try:
                sns.heatmap(corr, annot=True, cmap="coolwarm", ax=ax)
            except Exception:
                cax = ax.matshow(corr, cmap="coolwarm")
                fig.colorbar(cax)
                ax.set_xticks(range(len(corr.columns)))
                ax.set_xticklabels(corr.columns, rotation=90)
                ax.set_yticks(range(len(corr.columns)))
                ax.set_yticklabels(corr.columns)
            st.pyplot(fig)

        # Flexible Top/Bottom-N (any column)
        st.markdown("---")
        st.markdown("### Pilih kolom untuk Top/Bottom-N")
        col_list = list(df.columns)
        if col_list:
            widget_prefix = f"{filename}___{sheet_name}"
            chosen_col = st.selectbox("Pilih kolom:", options=col_list, key=f"col_{widget_prefix}")
            max_n = min(100, max(1, len(df)))
            chosen_n = st.slider("Jumlah baris (N):", min_value=1, max_value=max_n, value=min(10, max_n), key=f"n_{widget_prefix}")
            sort_order = st.radio("Urutan:", ["Top N (descending)", "Bottom N (ascending)"], index=0, key=f"ord_{widget_prefix}", horizontal=True)
            if st.button("Tampilkan Top/Bottom N", key=f"btn_{widget_prefix}"):
                try:
                    if pd.api.types.is_numeric_dtype(df[chosen_col]):
                        asc = (sort_order == "Bottom N (ascending)")
                        st.dataframe(df.sort_values(by=chosen_col, ascending=asc).head(chosen_n))
                    else:
                        if sort_order == "Top N (descending)":
                            vc = df[chosen_col].value_counts().head(chosen_n)
                        else:
                            vc = df[chosen_col].value_counts(ascending=True).head(chosen_n)
                        st.dataframe(vc.to_frame(name="count"))
                except Exception as e:
                    st.warning(f"⚠️ Gagal menampilkan Top/Bottom-N untuk kolom {chosen_col}: {e}")

    # Download buttons (Excel & HTML)
    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            "⬇️ Download laporan Excel",
            data=out_excel,
            file_name=f"analysis_{filename}_{sheet_name}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    with col2:
        st.download_button(
            "⬇️ Download laporan HTML",
            data=html_bytes,
            file_name=f"analysis_{filename}_{sheet_name}.html",
            mime="text/html"
        )

# -------------------------
# UI (main) — keep original flows for non-tabular files unchanged
# -------------------------
st.title("🤖 Gemini 2.5 Flash + Groq Chatbot — Multi-files + Excel Analysis")
st.write("Upload banyak file (PDF, TXT, DOCX, PPTX, Images, CSV, XLS, XLSX). Gambar akan diproses dengan OCR.Space API. Excel/CSV akan dianalisa otomatis.")

# Sidebar upload & actions
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files (pdf, txt, docx, pptx, images, csv, xls, xlsx) — boleh banyak",
    type=["pdf", "txt", "docx", "pptx", "jpg", "jpeg", "png", "gif", "bmp", "jfif", "csv", "xls", "xlsx"],
    accept_multiple_files=True
)
build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

# parse uploaded files immediately to fill dataframes (but don't rebuild vector store until button)
if uploaded_files:
    for f in uploaded_files:
        if f.name not in st.session_state.dataframes:
            try:
                extract_text_from_file(f)
            except Exception as e:
                st.sidebar.warning(f"Failed parse {f.name}: {e}")

# Build vector store
if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file terlebih dahulu.")
    else:
        with st.spinner("📦 Memproses file dan membuat vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error("Tidak ada teks valid berhasil diekstrak. Periksa file.")
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(f"Vector store terbangun. Dokumen: {len(st.session_state.indexed_files)} | Chunk total: {len(docs)}")

# Reset
if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.session_state.dataframes = {}
    st.success("Vector store & dataframes di-reset.")

# Show indexed files
if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# Data preview & analysis for any parsed DataFrames
if st.session_state.dataframes:
    st.subheader("📊 Data Preview, Profiling & Analisa Otomatis (Excel/CSV)")
    show_in_app = st.checkbox("Tampilkan analisa di Streamlit (per sheet)", value=True)
    for fname, payload in st.session_state.dataframes.items():
        with st.expander(f"🔎 File: {fname}", expanded=False):
            for sheet_name, df in payload["sheets"].items():
                auto_analyze_dataframe(df, fname, sheet_name, show_in_app)

# -------------------------
# Prompting / Q&A area (keep original)
# -------------------------
model_choice = st.sidebar.radio(
    "Pilih LLM Provider:",
    ["Gemini 2.5 Flash (Google)", "Groq (llama-3.3-70b-versatile)"]
)

prompt = st.text_input(
    "Tanyakan sesuatu berdasarkan dokumen yang diupload:",
    placeholder="Misal: Ringkas dokumen tentang topik X..."
)
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload file dan klik 'Build Vector Store'.")
    else:
        with st.spinner("🔎 Mengambil konteks dari vector store..."):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)
        # format & system instructions (kept simple)
        context_text = "\n\n".join([d.page_content for d in results])
        composed_prompt = (
            "Jawablah seakurat mungkin berdasarkan konteks berikut.\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYAAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )
        try:
            if model_choice.startswith("Gemini"):
                from langchain_google_genai import ChatGoogleGenerativeAI
                llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2)
                with st.spinner("🤖 Gemini sedang menjawab..."):
                    response = llm.invoke(composed_prompt)
            else:
                from langchain_groq import ChatGroq
                llm = ChatGroq(temperature=0.2, groq_api_key=GROQ_API_KEY, model_name="llama-3.3-70b-versatile")
                with st.spinner("⚡ Groq sedang menjawab..."):
                    response = llm.invoke(composed_prompt)
            st.subheader("💬 Jawaban")
            out_text = getattr(response, "content", None) or str(response)
            st.write(out_text)
        except Exception as e:
            st.error(f"❌ Error saat memanggil LLM: {e}")
